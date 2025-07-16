import os
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.documents import Document
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.chat_models import ChatOllama

def load_pptx_text(pptx_path):
    """
    Extracts and returns all text from a .pptx file as a list of Documents.
    """
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        if slide_text:
            full_text = "\n".join(slide_text).strip()
            texts.append(Document(page_content=full_text))
    return texts

def summarize_pptx():
    """
    Main function to summarize PPTX using Ollama models locally.
    """
    pptx_path = input("Please enter the full path to your PPTX file: ")
    if not os.path.exists(pptx_path):
        print(f"Error: The file '{pptx_path}' does not exist.")
        return

    try:
        # --- Load and Split ---
        print("Loading and splitting PPTX...")
        docs = load_pptx_text(pptx_path)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = text_splitter.split_documents(docs)

        # --- Embeddings using Ollama ---
        print("Creating vector store with Ollama embeddings...")
        embeddings = OllamaEmbeddings(model="snowflake-arctic-embed") 
        vectorstore = FAISS.from_documents(splits, embeddings)
        retriever = vectorstore.as_retriever()

        # --- Ollama LLM for Summarization ---
        print("Loading Ollama LLM...")
        llm = ChatOllama(model="koesn/llama3-8b-instruct:latest")  

        # --- Prompt Template ---
        prompt_template = """Use the following context to answer the question:
        {context}

        Question: {question}
        """
        prompt = ChatPromptTemplate.from_template(prompt_template)

        def format_docs(docs):
            return "\n\n".join(doc.page_content for doc in docs)

        rag_chain = (
            {"context": retriever | format_docs, "question": RunnablePassthrough()}
            | prompt
            | llm
            | StrOutputParser()
        )

        # --- Interactive Loop ---
        print("\nPPTX processed. You can now ask questions.")
        while True:
            query = input("Your question (or type 'exit' to quit): ")
            if query.lower() == 'exit':
                break
            print("Generating summary...")
            result = rag_chain.invoke(query)
            print("\n--- Summary ---\n")
            print(result)
            print("\n---------------\n")

    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    summarize_pptx()
